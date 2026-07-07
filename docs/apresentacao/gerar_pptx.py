# -*- coding: utf-8 -*-
"""Gera SunTrip_Apresentacao.pptx na pasta docs/apresentacao/"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BASE = Path(__file__).resolve().parent
IMG = BASE / "imagens"
OUT = BASE / "SunTrip_Apresentacao.pptx"

NAVY = RGBColor(12, 26, 46)
OCEAN = RGBColor(14, 165, 233)
YELLOW = RGBColor(250, 204, 21)
WHITE = RGBColor(241, 245, 249)
MUTED = RGBColor(148, 163, 184)


def set_slide_bg(slide, prs):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY


def add_title_box(slide, title, subtitle=None, top=0.8):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(8.8), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = MUTED
        p2.space_before = Pt(8)


def add_bullets(slide, items, top=2.0, font_size=20):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(8.6), Inches(4.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = WHITE
        p.space_after = Pt(10)


def add_image_slide(prs, title, image_name, caption=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, prs)
    add_title_box(slide, title, caption, top=0.4)
    path = IMG / image_name
    if path.exists():
        slide.shapes.add_picture(
            str(path),
            Inches(0.8),
            Inches(1.35),
            width=Inches(8.4),
        )
    else:
        add_bullets(slide, [f"(Imagem não encontrada: {image_name})"], top=2.5)


def add_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, prs)
    # faixa amarela
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(3.2), Inches(10), Inches(0.08)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = YELLOW
    shape.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "SunTrip"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = YELLOW
    p2 = tf.add_paragraph()
    p2.text = "Pagamentos digitais para transporte em Angola"
    p2.font.size = Pt(24)
    p2.font.color.rgb = WHITE
    p2.space_before = Pt(12)
    p3 = tf.add_paragraph()
    p3.text = "Fintech · Carteira · QR Code · Kwanza (Kz)"
    p3.font.size = Pt(16)
    p3.font.color.rgb = MUTED
    p3.space_before = Pt(20)
    p4 = tf.add_paragraph()
    p4.text = "Apresentação do projecto · 2026"
    p4.font.size = Pt(14)
    p4.font.color.rgb = OCEAN
    p4.space_before = Pt(36)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1 Capa
    add_cover(prs)

    # 2 Introdução
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, prs)
    add_title_box(slide, "Introdução", "O que é o SunTrip?")
    add_bullets(
        slide,
        [
            "Plataforma web de pagamentos digitais para táxis e transporte urbano.",
            "Permite pagar corridas, gerir carteira e receber via QR Code — tudo em Kwanza (Kz).",
            "Dois perfis: Cliente (passageiro) e Motorista (taxista).",
            "Objectivo: modernizar pagamentos no transporte angolano de forma segura e simples.",
        ],
    )

    # 3 Importância
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, prs)
    add_title_box(slide, "Importância", "Por que este projecto importa?")
    add_bullets(
        slide,
        [
            "Reduz dependência de dinheiro físico nas corridas.",
            "Registo digital com validação (email, telefone angolano, dados bancários).",
            "Histórico transparente: data, hora, valor, estado e motorista/cliente.",
            "Adaptado à realidade local: domínios de email comuns, telefone com 9 dígitos.",
            "Base para escalar pagamentos digitais no sector do transporte.",
        ],
    )

    # 4 Funcionalidades
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, prs)
    add_title_box(slide, "Funcionalidades principais")
    add_bullets(
        slide,
        [
            "Autenticação: registo Cliente/Motorista, login, logout (JWT).",
            "Carteira virtual: adicionar saldo, transferir, consultar saldo.",
            "QR Code: motorista gera QR · cliente paga com a carteira.",
            "Transferência bancária (IBAN / conta).",
            "Histórico com filtros (hoje, semana, mês) e estados: Pago, Pendente, Cancelado.",
            "Dashboards separados + painel administrativo.",
        ],
        font_size=18,
    )

    # 5-9 Imagens
    add_image_slide(prs, "Ecrã de login", "01-login.png", "Acesso seguro com email validado")
    add_image_slide(prs, "Escolha do tipo de conta", "02-registo-tipo.png", "Cliente ou Taxista / Motorista")
    add_image_slide(prs, "Registo de Cliente", "03-registo-cliente.png", "Nome, telefone, email, conta bancária")
    add_image_slide(prs, "Registo de Motorista", "04-registo-motorista.png", "Carta, matrícula, documento de identificação")
    add_image_slide(
        prs,
        "Dashboard do Motorista",
        "05-dashboard-motorista.png",
        "Saldo, QR Code, últimos pagamentos recebidos",
    )

    # 10 Tecnologias
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, prs)
    add_title_box(slide, "Tecnologias e linguagens")
    add_bullets(
        slide,
        [
            "Frontend: JavaScript · React · Next.js 14 · Tailwind CSS",
            "Backend: JavaScript · Node.js · Express.js",
            "Base de dados: MongoDB · Mongoose",
            "Segurança: JWT · bcrypt (palavras-passe)",
            "Tempo real: Socket.io",
            "Arquitectura: API REST (porta 5000) + Web App (porta 3000)",
        ],
        font_size=18,
    )

    # 11 Desenvolvimento
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, prs)
    add_title_box(slide, "Desenvolvimento", "Como foi construído")
    add_bullets(
        slide,
        [
            "Estrutura em monorepo: client/ (interface) + server/ (API).",
            "Validações no frontend e backend (email, telefone 9XXXXXXXX).",
            "Contas começam com saldo 0 Kz — dinheiro só entra pela app.",
            "Testes em PC e telemóvel na mesma rede Wi‑Fi.",
            "Scripts de arranque: INICIAR-SUNTRIP.bat · seed de contas demo.",
        ],
        font_size=18,
    )

    # 12 Conclusão
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, prs)
    add_title_box(slide, "Conclusão")
    add_bullets(
        slide,
        [
            "O SunTrip demonstra uma solução completa de fintech para transporte em Angola.",
            "Interface moderna, responsiva e pensada para clientes e motoristas.",
            "Projecto funcional: registo, carteira, QR, histórico e dashboards.",
            "Próximos passos possíveis: app móvel nativa, notificações, integração bancária real.",
            "",
            "Demo: cliente@gmail.com / motorista@gmail.com · palavra-passe: demo123",
            "Obrigado pela atenção!",
        ],
        font_size=19,
    )

    prs.save(str(OUT))
    print(f"Apresentação criada: {OUT}")


if __name__ == "__main__":
    main()
